import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

THEMES = {
    "professional": {
        "bg": (15, 15, 25),
        "title_fg": (255, 255, 255),
        "body_fg": (220, 220, 230),
        "accent": (0, 180, 130),
        "bullet_fg": (200, 200, 215),
        "notes_fg": (150, 150, 160),
    },
    "minimal": {
        "bg": (250, 250, 250),
        "title_fg": (20, 20, 30),
        "body_fg": (50, 50, 60),
        "accent": (30, 100, 220),
        "bullet_fg": (60, 60, 80),
        "notes_fg": (120, 120, 140),
    },
    "vibrant": {
        "bg": (20, 10, 50),
        "title_fg": (255, 255, 255),
        "body_fg": (220, 210, 240),
        "accent": (255, 80, 140),
        "bullet_fg": (200, 190, 220),
        "notes_fg": (160, 150, 180),
    },
}


def _rgb(t):
    return RGBColor(*t)


class PPTTool:
    def __init__(self, output_dir: Path):
        self.output_dir = output_dir

    def create_ppt(self, slide_data: dict, style: str = "professional"):
        theme = THEMES.get(style, THEMES["professional"])
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # --- Title slide ---
        title_slide = prs.slides.add_slide(blank_layout)
        self._add_bg(title_slide, theme["bg"])

        # Centered title box
        txb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "Presentation")
        p.runs[0].font.size = Pt(40)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = _rgb(theme["title_fg"])
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER

        # Date subtitle
        date_txb = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(0.5))
        dtf = date_txb.text_frame
        dp = dtf.paragraphs[0]
        dp.text = datetime.now().strftime("%B %d, %Y")
        dp.runs[0].font.size = Pt(16)
        dp.runs[0].font.color.rgb = _rgb(theme["accent"])
        dp.alignment = PP_ALIGN.CENTER

        # Accent line
        from pptx.util import Emu
        line = title_slide.shapes.add_shape(
            1, Inches(4.5), Inches(3.9), Inches(4.33), Emu(36000)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(theme["accent"])
        line.line.fill.background()

        # --- Content slides ---
        for i, slide_info in enumerate(slide_data.get("slides", []), 1):
            slide = prs.slides.add_slide(blank_layout)
            self._add_bg(slide, theme["bg"])

            # Header bar
            header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
            header.fill.solid()
            header.fill.fore_color.rgb = _rgb(tuple(max(0, c - 8) for c in theme["bg"]))
            header.line.fill.background()

            # Accent sidebar
            sidebar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(7.5))
            sidebar.fill.solid()
            sidebar.fill.fore_color.rgb = _rgb(theme["accent"])
            sidebar.line.fill.background()

            # Slide number
            num_txb = slide.shapes.add_textbox(Inches(12.5), Inches(0.1), Inches(0.7), Inches(0.5))
            ntf = num_txb.text_frame
            np_ = ntf.paragraphs[0]
            np_.text = str(i)
            np_.runs[0].font.size = Pt(11)
            np_.runs[0].font.color.rgb = _rgb(theme["accent"])

            # Title
            title_txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12), Inches(0.85))
            ttf = title_txb.text_frame
            tp = ttf.paragraphs[0]
            tp.text = slide_info.get("title", "")
            tp.runs[0].font.size = Pt(24)
            tp.runs[0].font.bold = True
            tp.runs[0].font.color.rgb = _rgb(theme["title_fg"])

            # Separator line
            sep = slide.shapes.add_shape(1, Inches(0.3), Inches(1.15), Inches(12.7), Emu(18000))
            sep.fill.solid()
            sep.fill.fore_color.rgb = _rgb(theme["accent"])
            sep.line.fill.background()

            # Bullet points
            bullets = slide_info.get("bullet_points", [])
            if bullets:
                body_txb = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8))
                btf = body_txb.text_frame
                btf.word_wrap = True
                for j, bullet in enumerate(bullets):
                    if j == 0:
                        bp = btf.paragraphs[0]
                    else:
                        bp = btf.add_paragraph()
                    bp.text = f"▸  {bullet}"
                    if bp.runs:
                        bp.runs[0].font.size = Pt(17)
                        bp.runs[0].font.color.rgb = _rgb(theme["bullet_fg"])
                    from pptx.util import Pt as Pt2
                    bp.space_before = Pt2(6)

            # Speaker notes
            notes_text = slide_info.get("speaker_notes", "")
            if notes_text:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = notes_text

        # Save
        raw_title = slide_data.get("title", "Presentation")
        safe_title = re.sub(r"[^\w\s-]", "", raw_title).strip().replace(" ", "_")[:40]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_title}_{timestamp}.pptx"
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        return filename, str(filepath)

    def _add_bg(self, slide, color_tuple):
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(color_tuple)
        bg.line.fill.background()
